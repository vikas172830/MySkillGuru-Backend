import asyncio
import csv
import io
import logging
import os
import re
import tempfile
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Tuple

import requests
from bson import ObjectId
from google import genai
from google.genai import types as genai_types
from motor.motor_asyncio import AsyncIOMotorDatabase

from app.core.config import settings
from app.utils.token_usage import increment_institute_gemini_tokens

_client: genai.Client | None = None


def _get_client() -> genai.Client:
    """Lazy singleton — avoids failing app startup when GEMINI_API_KEY isn't set yet."""
    global _client
    if _client is None:
        _client = genai.Client(api_key=settings.GEMINI_API_KEY)
    return _client

try:
    from docx import Document as DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False
    logging.warning("python-docx not installed — .docx/.doc extraction unavailable")

try:
    from pptx import Presentation as PptxPresentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False
    logging.warning("python-pptx not installed — .pptx/.ppt extraction unavailable")

try:
    import openpyxl
    _HAS_XLSX = True
except ImportError:
    _HAS_XLSX = False
    logging.warning("openpyxl not installed — .xlsx/.xls extraction unavailable")


# ==========================================
# FILE-TYPE HELPERS
# ==========================================

def _file_ext(filename: str) -> str:
    return Path(filename).suffix.lstrip(".").lower()


def _extract_text_from_docx(file_bytes: bytes) -> str:
    if not _HAS_DOCX:
        return ""
    doc = DocxDocument(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text.strip())
    return "\n".join(parts)


def _extract_text_from_doc(file_bytes: bytes) -> str:
    """Legacy .doc: try python-docx, then LibreOffice headless, then raw byte strip."""
    try:
        if _HAS_DOCX:
            return _extract_text_from_docx(file_bytes)
    except Exception:
        pass

    try:
        with tempfile.TemporaryDirectory() as tmp_dir:
            doc_path = os.path.join(tmp_dir, "input.doc")
            with open(doc_path, "wb") as f:
                f.write(file_bytes)
            ret = os.system(
                f"libreoffice --headless --convert-to txt:Text --outdir {tmp_dir} {doc_path} > /dev/null 2>&1"
            )
            txt_path = os.path.join(tmp_dir, "input.txt")
            if ret == 0 and os.path.exists(txt_path):
                with open(txt_path, "r", errors="replace") as f:
                    return f.read()
    except Exception as e:
        logging.warning("LibreOffice .doc fallback failed: %s", e)

    raw = file_bytes.decode("latin-1", errors="replace")
    printable = "".join(c if c.isprintable() or c in "\n\t " else " " for c in raw)
    return re.sub(r"[ \t]{3,}", " ", re.sub(r"\n{3,}", "\n\n", printable)).strip()


def _extract_text_from_pptx(file_bytes: bytes) -> str:
    if not _HAS_PPTX:
        return ""
    prs = PptxPresentation(io.BytesIO(file_bytes))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def _extract_text_from_xlsx(file_bytes: bytes) -> str:
    if not _HAS_XLSX:
        return ""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(c) for c in row if c is not None)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def _extract_text_from_csv(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    return "\n".join("\t".join(row) for row in reader)


def _extract_text_from_plain(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="replace")


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Route file_bytes to the correct local extractor based on extension. PDFs go straight to Gemini."""
    ext = _file_ext(filename)

    extractors = {
        "docx": _extract_text_from_docx,
        "doc": _extract_text_from_doc,
        "pptx": _extract_text_from_pptx,
        "ppt": _extract_text_from_pptx,
        "xlsx": _extract_text_from_xlsx,
        "xls": _extract_text_from_xlsx,
        "txt": _extract_text_from_plain,
        "md": _extract_text_from_plain,
        "csv": _extract_text_from_csv,
        "odt": _extract_text_from_plain,
    }

    fn = extractors.get(ext)
    if fn is None:
        return file_bytes.decode("utf-8", errors="replace")

    try:
        return fn(file_bytes)
    except Exception as e:
        logging.error("extract_text_from_file(%s) error: %s", ext, e)
        return ""


# ==========================================
# GEMINI
# ==========================================

_GEMINI_PDF_PROMPT = (
    "Extract ALL text from this question paper PDF with high accuracy. "
    "Preserve the original structure including:\n"
    "- Section headings\n"
    "- Question numbers\n"
    "- Sub questions (a/b/c)\n"
    "- Marks allocation\n"
    "- Mathematical expressions\n"
    "- Tables\n"
    "- Chemical formulas\n\n"
    "Return plain text only.\n"
    "Do not use markdown."
)

_GEMINI_TEXT_PROMPT = (
    "You are given text extracted from a question paper. "
    "Reformat it accurately, preserving:\n"
    "- Section headings\n"
    "- Question numbers and sub-questions (a/b/c)\n"
    "- Marks allocation\n"
    "- Mathematical expressions\n"
    "- Tables and chemical formulas\n\n"
    "Return plain text only. Do not use markdown."
)


def _parse_token_usage(usage_metadata) -> dict:
    try:
        return {
            "prompt_tokens": getattr(usage_metadata, "prompt_token_count", 0) or 0,
            "candidate_tokens": getattr(usage_metadata, "candidates_token_count", 0) or 0,
            "total_tokens": getattr(usage_metadata, "total_token_count", 0) or 0,
        }
    except Exception:
        return {"prompt_tokens": 0, "candidate_tokens": 0, "total_tokens": 0}


def extract_question_paper_text_with_gemini(file_bytes: bytes, filename: str = "file.pdf") -> Tuple[str, dict]:
    """
    PDF -> sent as binary part directly (native Gemini vision).
    Other formats -> extracted locally, then sent to Gemini as text for clean reformatting.
    Blocking call — run via asyncio.to_thread() from async callers.
    """
    ext = _file_ext(filename)
    client = _get_client()

    if ext == "pdf":
        pdf_part = genai_types.Part.from_bytes(data=file_bytes, mime_type="application/pdf")
        response = client.models.generate_content(model="gemini-2.5-flash", contents=[pdf_part, _GEMINI_PDF_PROMPT])
        text = ""
        try:
            text = response.text.strip()
        except Exception:
            pass
        return text, _parse_token_usage(getattr(response, "usage_metadata", None))

    raw_text = extract_text_from_file(file_bytes, filename)
    if not raw_text.strip():
        logging.warning("[qp-extract] No text extracted locally from '%s'", filename)
        return "", {"prompt_tokens": 0, "candidate_tokens": 0, "total_tokens": 0}

    user_content = f"{_GEMINI_TEXT_PROMPT}\n\n---\n\n{raw_text}"
    response = client.models.generate_content(model="gemini-2.5-flash", contents=[user_content])
    text = ""
    try:
        text = response.text.strip()
    except Exception:
        pass

    return text, _parse_token_usage(getattr(response, "usage_metadata", None))


def _call_generate_content_with_retry(client: genai.Client, model: str, contents: list, config=None):
    """
    Shared retry/backoff for Gemini generate_content calls: up to 5 attempts on
    transient errors (503/UNAVAILABLE/429/RESOURCE_EXHAUSTED/quota), honoring a
    "Please retry in Ns" hint from the API when present. Returns the raw response
    object (not just text) so callers that need finish_reason/candidates — e.g.
    pdf_extract.py's truncation detection — can inspect it directly.

    Blocking — call via asyncio.to_thread() from async callers.
    """
    response = None
    for attempt in range(5):
        try:
            response = client.models.generate_content(model=model, contents=contents, config=config)
            break
        except Exception as e:
            error_text = str(e)
            is_retryable = (
                "503" in error_text
                or "UNAVAILABLE" in error_text
                or "429" in error_text
                or "RESOURCE_EXHAUSTED" in error_text
                or "quota" in error_text.lower()
            )
            if not is_retryable:
                raise RuntimeError(f"File extraction failed: {e}") from e

            match = re.search(r"Please retry in ([\d\.]+)s", error_text)
            if match:
                wait_time = float(match.group(1)) + 1.0
            elif "429" in error_text or "RESOURCE_EXHAUSTED" in error_text:
                wait_time = 5 * (2 ** attempt)
            else:
                wait_time = 2 ** attempt

            logging.warning(
                "Gemini request failed (retryable error). Retry %d/5 in %.2fs.", attempt + 1, wait_time
            )
            time.sleep(wait_time)

    if response is None:
        raise RuntimeError("Gemini service unavailable after 5 retries.")

    return response


def _response_hit_token_cap(response) -> bool:
    """
    True when Gemini stopped because it ran into the model's output-token
    ceiling, meaning the text returned is a fragment rather than the whole
    extraction. Without this check an oversized document comes back
    truncated but otherwise indistinguishable from a complete, successful
    result, and everything downstream treats the fragment as the full file.
    """
    try:
        finish_reason = response.candidates[0].finish_reason
        return finish_reason is not None and "MAX_TOKENS" in str(finish_reason)
    except Exception:
        return False


def generate_content_from_file_checked(
    file_bytes: bytes, mime_type: str, prompt: str, model: str = "gemini-2.5-flash"
) -> Tuple[str, dict, bool]:
    """
    Same as generate_content_from_file, plus a third return value reporting
    whether the response was cut short by the model's output-token cap.

    Callers extracting text out of a whole file need that signal — see
    app/api/routers/ai_tutor.py, which surfaces it to the student instead of
    silently building notes from a partial document.

    Blocking call — run via asyncio.to_thread() from async callers.
    """
    client = _get_client()
    file_part = genai_types.Part.from_bytes(data=file_bytes, mime_type=mime_type)
    response = _call_generate_content_with_retry(client, model, [file_part, prompt])

    text = ""
    try:
        text = response.text.strip()
    except Exception:
        pass

    usage = _parse_token_usage(getattr(response, "usage_metadata", None))
    return text, usage, _response_hit_token_cap(response)


def generate_content_from_file(
    file_bytes: bytes, mime_type: str, prompt: str, model: str = "gemini-2.5-flash"
) -> Tuple[str, dict]:
    """
    Send file bytes (PDF/image) directly to Gemini with a prompt. Blocking call —
    run via asyncio.to_thread() from async callers. Shared by both the
    homework-help and notes-generation pipelines (Flask only had this retry logic on the
    homework path; consolidating makes notes-generation equally robust for free).

    Truncation-unaware by design, for callers that don't act on the signal;
    use generate_content_from_file_checked when a partial result matters.
    """
    text, usage, _ = generate_content_from_file_checked(file_bytes, mime_type, prompt, model)
    return text, usage


def generate_html_from_prompt(
    prompt: str, model: str = "gemini-2.5-flash", max_tokens: int = 16000
) -> Tuple[str, dict]:
    """
    Text prompt -> complete HTML document. The Gemini counterpart to
    app.services.claude.generate_html, returning the same (html, usage)
    shape so call sites can swap providers without special-casing either.

    Used for generation that works from the model's own knowledge rather
    than from a long source document — see ai_tutor.py's notes job, which
    routes prompt-only requests here and keeps document-backed ones on
    Claude. Markdown fences are stripped and a truncation banner is appended
    on a short read, matching generate_html's behavior exactly, so the
    rendered PDF looks the same whichever provider produced it.

    Blocking call — run via asyncio.to_thread() from async callers.
    """
    client = _get_client()
    response = _call_generate_content_with_retry(
        client, model, [prompt],
        config=genai_types.GenerateContentConfig(
            max_output_tokens=max_tokens,
            # gemini-2.5-flash is a thinking model and charges its internal
            # reasoning to the SAME max_output_tokens budget as the visible
            # answer. Producing a formatted document from an explicit prompt
            # needs no chain-of-thought, so thinking is disabled to give the
            # whole budget to output — measured without this, thinking alone
            # consumed ~45% of a 3000-token budget before any HTML was
            # written, and the response came back cut off mid-<style>.
            thinking_config=genai_types.ThinkingConfig(thinking_budget=0),
        ),
    )

    # response.text is None (not "") when the model returned no text part —
    # re.sub would raise TypeError on it.
    html_content = ""
    try:
        html_content = response.text or ""
    except Exception:
        pass

    html_content = re.sub(r"^```html\s*", "", html_content, flags=re.IGNORECASE)
    html_content = re.sub(r"^```\s*", "", html_content)
    html_content = re.sub(r"\s*```\s*$", "", html_content)
    html_content = html_content.strip()

    if _response_hit_token_cap(response):
        logging.warning(
            "generate_html_from_prompt: output truncated at max_output_tokens=%d (model=%s)",
            max_tokens, model,
        )
        html_content += (
            '\n<div style="margin:24px;padding:12px 16px;border:1px solid #f59e0b;'
            'background:#fffbeb;color:#92400e;border-radius:6px;font-family:Arial,sans-serif;'
            'font-size:13px">This document was cut off before it finished generating. '
            'Try a shorter length, or regenerate.</div>'
        )

    return html_content, _parse_token_usage(getattr(response, "usage_metadata", None))


# ==========================================
# BACKGROUND JOB (FastAPI BackgroundTasks, not a raw Thread)
# ==========================================

async def extract_and_patch_question_paper_text(
    db: AsyncIOMotorDatabase,
    folder_id: ObjectId,
    question_paper_url: str,
    faculty_id: str,
    filename: str = "file.pdf",
) -> None:
    try:
        logging.info("[qp-extract] Downloading file for folder %s", folder_id)

        resp = await asyncio.to_thread(requests.get, question_paper_url, timeout=60)
        resp.raise_for_status()
        file_bytes = resp.content

        logging.info("[qp-extract] Extracting text for folder %s (type=%s)", folder_id, _file_ext(filename))

        extracted_text, token_usage = await asyncio.to_thread(
            extract_question_paper_text_with_gemini, file_bytes, filename
        )

        now = datetime.now(timezone.utc)
        await db["newsavedDocs"].update_one(
            {"_id": folder_id},
            {"$set": {"question_paper.text": extracted_text, "question_paper.text_at": now, "updated_at": now}},
        )

        await increment_institute_gemini_tokens(
            db, faculty_id, token_usage["prompt_tokens"], token_usage["candidate_tokens"]
        )

        logging.info("[qp-extract] Done for folder %s", folder_id)

    except requests.RequestException as e:
        logging.error("[qp-extract] File download failed: %s", e)
        await db["newsavedDocs"].update_one(
            {"_id": folder_id},
            {"$set": {"question_paper.text_error": str(e), "updated_at": datetime.now(timezone.utc)}},
        )

    except Exception as e:
        logging.error("[qp-extract] Extraction failed: %s", e)
        await db["newsavedDocs"].update_one(
            {"_id": folder_id},
            {"$set": {"question_paper.text_error": str(e), "updated_at": datetime.now(timezone.utc)}},
        )
